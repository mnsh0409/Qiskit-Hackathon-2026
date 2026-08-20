import sys, subprocess, tempfile, os
from pptx import Presentation
from pptx.util import Emu

def build(pdf_path, pptx_path, dpi=150):
    tmpdir = tempfile.mkdtemp()
    prefix = os.path.join(tmpdir, "p")
    subprocess.run(["pdftoppm", "-r", str(dpi), "-png", pdf_path, prefix], check=True)
    pages = sorted(f for f in os.listdir(tmpdir) if f.endswith(".png"))
    assert pages, "no pages rendered"

    from PIL import Image
    im = Image.open(os.path.join(tmpdir, pages[0]))
    w_px, h_px = im.size
    emu_per_px = 914400 / dpi
    slide_w = Emu(int(w_px * emu_per_px))
    slide_h = Emu(int(h_px * emu_per_px))

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    blank = prs.slide_layouts[6]

    for pg in pages:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(os.path.join(tmpdir, pg), 0, 0, width=slide_w, height=slide_h)

    prs.save(pptx_path)
    print(f"{pptx_path}: {len(pages)} slides, {w_px}x{h_px}px @ {dpi}dpi")

if __name__ == "__main__":
    pdf_path, pptx_path = sys.argv[1], sys.argv[2]
    build(pdf_path, pptx_path)
